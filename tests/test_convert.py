#!/usr/bin/env python3
"""test_convert.py — hermetic tests for the conversion pipeline
(scripts/convert.py + scripts/convert_formats.py).

No docling, no network. Heavy-engine formats (PDF/DOCX) are exercised only at
the routing level; PPTX/HTML handler tests run when python-pptx/markdownify
are importable and are SKIPPED (not failed) otherwise. Everything else —
the numbering heuristic, EML parsing, long-path helper, and the batch
driver's bookkeeping (resume, flags, errors, policy skips) — is stdlib-only.

Usage:
  python3 tests/test_convert.py
"""
import importlib.util
import io
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SCRIPTS = ROOT / "scripts"

_passed = 0
_skipped = []


def assert_true(name, cond, hint=""):
    global _passed
    if not cond:
        print(f"FAIL {name}" + (f"  [{hint}]" if hint else ""))
        sys.exit(1)
    _passed += 1
    print(f"OK   {name}")


def assert_eq(name, expected, actual):
    assert_true(name, expected == actual,
                hint=f"expected {expected!r}, got {actual!r}")


def skip(name, why):
    _skipped.append(name)
    print(f"SKIP {name} ({why})")


def import_script(name, path):
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


fmts = import_script("convert_formats", SCRIPTS / "convert_formats.py")


# ─── routing tables ──────────────────────────────────────────────────────────
def test_routing_tables():
    for ext in (".pdf", ".docx", ".doc", ".pptx", ".html", ".htm", ".eml", ".msg"):
        assert_true(f"handler registered for {ext}", ext in fmts.HANDLERS)
    for ext in (".md", ".txt"):
        assert_true(f"{ext} passes through", ext in fmts.PASSTHROUGH)
    for ext in (".xlsx", ".xls", ".xlsm", ".csv"):
        assert_true(f"{ext} policy-skipped", ext in fmts.POLICY_SKIP)
    overlap = (set(fmts.HANDLERS) & fmts.PASSTHROUGH) | \
              (set(fmts.HANDLERS) & fmts.POLICY_SKIP) | \
              (fmts.PASSTHROUGH & fmts.POLICY_SKIP)
    assert_eq("routing tables are disjoint", set(), overlap)


# ─── helpers ─────────────────────────────────────────────────────────────────
def test_wp_long_path_prefix():
    p = fmts.wp("some/relative/path.md")
    if os.name == "nt":
        assert_true("wp() adds \\\\?\\ prefix on Windows",
                    p.startswith("\\\\?\\"), hint=p)
        assert_eq("wp() is idempotent", p, fmts.wp(p))
    else:
        assert_true("wp() returns absolute path", os.path.isabs(p))


def test_visible_len_ignores_whitespace():
    assert_eq("visible_len strips whitespace", 4,
              fmts._visible_len("a b\n\tc  d\r\n"))


# ─── DOCX numbering heuristic ────────────────────────────────────────────────
def test_numbering_consistent_text_not_flagged():
    text = "1. Intro\n1.1 Scope\n1.2 Terms\n2. Body\n2.1 Detail\n"
    assert_eq("consistent numbering raises no flag", [],
              fmts._numbering_check(text))


def test_numbering_fabricated_text_flagged():
    # Sub-numbers claim section 3/4 while the document is still in section 1 —
    # the signature of docling reconstructing Word auto-numbering wrongly.
    text = "1. Intro\n3.1 Ghost\n4.2 Ghost\n1.1 Real\n"
    flags = fmts._numbering_check(text)
    assert_eq("fabricated numbering raises one flag", 1, len(flags))
    assert_eq("flag type is 'numbering'", "numbering", flags[0]["type"])
    assert_true("flag reports the mismatch ratio", "2/3" in flags[0]["detail"],
                hint=flags[0]["detail"])


def test_numbering_bulleted_lines_counted():
    text = "1. Intro\n- 2.1 bulleted ghost\n"
    flags = fmts._numbering_check(text)
    assert_eq("bulleted sub-number still checked", 1, len(flags))


# ─── EML (stdlib) ────────────────────────────────────────────────────────────
EML_SAMPLE = b"""From: Alice Example <alice@example.org>
To: Bob Example <bob@example.org>
Subject: Quarterly report attached
Date: Mon, 01 Jan 2026 10:00:00 +0000
MIME-Version: 1.0
Content-Type: multipart/mixed; boundary="XYZ"

--XYZ
Content-Type: text/plain; charset="utf-8"

Please find the report attached.

> Earlier quoted reply here.
--XYZ
Content-Type: application/pdf; name="report.pdf"
Content-Disposition: attachment; filename="report.pdf"
Content-Transfer-Encoding: base64

JVBERi0=
--XYZ--
"""


def test_eml_frontmatter_body_and_attachments():
    with tempfile.TemporaryDirectory() as tmpdir:
        eml = Path(tmpdir) / "mail.eml"
        eml.write_bytes(EML_SAMPLE)
        text, flags = fmts.convert_eml(str(eml))
    assert_true("frontmatter opens the file", text.startswith("---\n"))
    assert_true("subject in frontmatter",
                'subject: "Quarterly report attached"' in text)
    assert_true("from in frontmatter", "alice@example.org" in text)
    assert_true("plain body preferred", "Please find the report attached." in text)
    assert_true("quoted chain kept", "> Earlier quoted reply here." in text)
    assert_true("attachment listed, not inlined",
                "**Attachments (not inlined):** report.pdf" in text)
    assert_eq("one attachments flag", 1, len(flags))
    assert_eq("flag type", "attachments", flags[0]["type"])


def test_eml_utf8_body_without_declared_charset():
    """Wild-caught case: UTF-8 body, no Content-Type charset. The stdlib
    default decode mangles it to replacement chars; the handler must fall
    back to a clean UTF-8 decode of the raw payload."""
    with tempfile.TemporaryDirectory() as tmpdir:
        eml = Path(tmpdir) / "nocharset.eml"
        eml.write_bytes(
            "From: a@b.c\nSubject: 한국어 제목\n\n본문 내용입니다.\n".encode("utf-8"))
        text, _ = fmts.convert_eml(str(eml))
    assert_true("undeclared-charset UTF-8 body survives",
                "본문 내용입니다" in text, hint=repr(text[-40:]))
    assert_true("no replacement characters in body", "�" not in text)


# ─── PPTX / HTML — engine-dependent, skip when the engine is absent ──────────
def test_pptx_slide_anchors_and_notes():
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        skip("pptx handler", "python-pptx not installed in this interpreter")
        return
    with tempfile.TemporaryDirectory() as tmpdir:
        deck = Path(tmpdir) / "deck.pptx"
        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[1])   # title + content
        s1.shapes.title.text = "First Slide"
        s1.placeholders[1].text = ("Body text long enough to stay above the "
                                   "low-text threshold. " * 8)
        s1.notes_slide.notes_text_frame.text = "Speaker note line."
        s2 = prs.slides.add_slide(prs.slide_layouts[6])   # blank → low text
        prs.save(str(deck))

        text, flags = fmts.convert_pptx(str(deck))
    assert_true("slide 1 anchor with title",
                "## Slide 1 — First Slide" in text, hint=text[:120])
    assert_true("slide 2 anchor present", "## Slide 2" in text)
    assert_true("speaker notes as blockquote", "> Speaker note line." in text)
    low = [f for f in flags if f["type"] == "pptx-low-text-slides"]
    assert_eq("one low-text flag", 1, len(low))
    assert_true("flag names slide 2 only", "[2]" in low[0]["detail"],
                hint=low[0]["detail"])


def test_html_conversion():
    try:
        import markdownify  # noqa: F401
        have_engine = True
    except ImportError:
        try:
            import html2text  # noqa: F401
            have_engine = True
        except ImportError:
            have_engine = False
    if not have_engine:
        skip("html handler", "neither markdownify nor html2text installed")
        return
    with tempfile.TemporaryDirectory() as tmpdir:
        page = Path(tmpdir) / "page.html"
        page.write_text("<h1>Title</h1><p>Body <b>bold</b>.</p>",
                        encoding="utf-8")
        text, flags = fmts.convert_html(str(page))
    assert_true("heading converted", "Title" in text)
    assert_true("bold body converted", "bold" in text)
    assert_eq("html raises no flags", [], flags)


# ─── batch driver bookkeeping (stdlib-only source tree) ──────────────────────
def _run_batch(src, dst, extra=()):
    env = dict(os.environ, PYTHONUTF8="1")
    return subprocess.run(
        [sys.executable, str(SCRIPTS / "convert.py"), "--batch",
         str(src), str(dst), "--workers", "1", *extra],
        capture_output=True, text=True, timeout=60, env=env)


def test_batch_bookkeeping_and_resume():
    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src"
        dst = Path(tmpdir) / "dst"
        (src / "Tier 1").mkdir(parents=True)
        (src / "Tier 1" / "note.md").write_text("# hi\n", encoding="utf-8")
        (src / "Tier 1" / "plain.txt").write_text("text\n", encoding="utf-8")
        (src / "Tier 1" / "data.xlsx").write_bytes(b"not really excel")
        (src / "Tier 1" / "image.png").write_bytes(b"\x89PNG")

        r = _run_batch(src, dst)
        assert_eq("batch exits 0 with no errors", 0, r.returncode)
        assert_true("md copied through",
                    (dst / "Tier 1" / "note.md").is_file())
        assert_true("txt copied through",
                    (dst / "Tier 1" / "plain.txt").is_file())
        assert_true("xlsx NOT converted",
                    not (dst / "Tier 1" / "data.xlsx.md").exists())
        skipped = (dst / ".convert_skipped.txt").read_text(encoding="utf-8")
        assert_true("xlsx recorded as excel-policy skip",
                    "excel-policy" in skipped, hint=skipped)
        assert_true("png recorded as no-handler skip",
                    "no handler for .png" in skipped, hint=skipped)
        done = (dst / ".convert_done.txt").read_text(encoding="utf-8")
        assert_eq("all 4 items marked done", 4,
                  len([l for l in done.splitlines() if l.strip()]))

        # Resume: nothing left to do, done-file unchanged.
        r2 = _run_batch(src, dst)
        assert_eq("resume exits 0", 0, r2.returncode)
        assert_true("resume processes 0 items",
                    "0 to process" in r2.stdout, hint=r2.stdout[-300:])


def test_batch_error_recorded_and_not_retried():
    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src"
        dst = Path(tmpdir) / "dst"
        src.mkdir()
        # Garbage bytes with a .pptx extension raise in the handler whether
        # python-pptx is installed (BadZipFile) or not (ImportError).
        (src / "broken.pptx").write_bytes(b"this is not a zip archive")
        r = _run_batch(src, dst)
        assert_eq("batch with an error exits 1", 1, r.returncode)
        errors = (dst / ".convert_errors.txt").read_text(encoding="utf-8")
        assert_true("error recorded with type", "broken.pptx" in errors,
                    hint=errors)
        done = (dst / ".convert_done.txt").read_text(encoding="utf-8")
        assert_true("errored file marked done (no resume stall)",
                    "broken.pptx" in done)
        r2 = _run_batch(src, dst)
        assert_true("errored file not retried on resume",
                    "0 to process" in r2.stdout, hint=r2.stdout[-300:])


def test_single_file_mode_stdout_and_flags():
    with tempfile.TemporaryDirectory() as tmpdir:
        eml = Path(tmpdir) / "mail.eml"
        eml.write_bytes(EML_SAMPLE)
        env = dict(os.environ, PYTHONUTF8="1")
        r = subprocess.run(
            [sys.executable, str(SCRIPTS / "convert.py"), str(eml)],
            capture_output=True, text=True, timeout=60, env=env)
        assert_eq("single-file mode exits 0", 0, r.returncode)
        assert_true("markdown on stdout", "Quarterly report attached" in r.stdout)
        assert_true("flags on stderr", "FLAG [attachments]" in r.stderr,
                    hint=r.stderr[:200])


def main():
    print("=== test_convert.py ===")
    test_routing_tables()
    test_wp_long_path_prefix()
    test_visible_len_ignores_whitespace()
    test_numbering_consistent_text_not_flagged()
    test_numbering_fabricated_text_flagged()
    test_numbering_bulleted_lines_counted()
    test_eml_frontmatter_body_and_attachments()
    test_eml_utf8_body_without_declared_charset()
    test_pptx_slide_anchors_and_notes()
    test_html_conversion()
    test_batch_bookkeeping_and_resume()
    test_batch_error_recorded_and_not_retried()
    test_single_file_mode_stdout_and_flags()
    tail = f" ({len(_skipped)} skipped: {', '.join(_skipped)})" if _skipped else ""
    print(f"\nAll conversion tests passed ({_passed} assertions){tail}.")


if __name__ == "__main__":
    main()
