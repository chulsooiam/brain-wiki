#!/usr/bin/env python3
"""Per-format -> Markdown converters for the corpus conversion pipeline.

Each handler takes a source path and returns (markdown_text, flags) where
flags is a list of {"type": ..., "detail": ...} dicts describing follow-up
work (vision layer, numbering review, attachments to route). Handlers never
write files; the driver (convert.py) owns all output.

Engine choices (measured on a ~400-document professional corpus, 2026-07/08):
- PDF / DOCX: docling, placeholder image mode (markitdown corrupts two-column
  PDFs and leaks Word hyperlink targets).
- PPTX: python-pptx, NOT docling — docling flattens decks into one stream
  with no slide boundaries (verified on a 16-slide deck: 0 headings), which
  destroys slide-number citation anchors. python-pptx keeps per-slide
  structure and is the only route to speaker notes.
- HTML: plain html->md library; already markup, docling adds nothing.
- EML/MSG: header-aware parsing, not layout conversion.
- XLSX and friends are policy-routed (catalogue, don't convert) by the driver.
"""
import glob
import io
import os
import re
import shutil
import subprocess
import threading

MIN_CHARS = 200        # below this a whole-document conversion is "low text"
LOWTEXT_SLIDE = 200    # per-slide threshold for the PPTX vision-layer flag
RECALL_MIN = 0.97      # docling vocabulary recall below this triggers recovery
RECALL_REPLACE = 0.70  # below this docling's structure isn't worth keeping


def wp(path):
    """Windows long-path safe form (bypasses the legacy 260-char MAX_PATH)."""
    p = os.path.abspath(path)
    if os.name == "nt" and not p.startswith("\\\\?\\"):
        return "\\\\?\\" + p
    return p


def _visible_len(text):
    return len("".join(ch for ch in text if not ch.isspace()))


# --------------------------------------------------------------------------
# PDF / DOCX via docling
# --------------------------------------------------------------------------
_tls = threading.local()


def _docling_converter():
    # One converter per thread: DocumentConverter is not documented as
    # thread-safe, and construction is cheap next to conversion.
    #
    # OCR is OFF by default (CONVERT_OCR=1 re-enables): the OCR preprocessor
    # rasterizes every page at high resolution, and image-heavy PDFs (maps,
    # dashboards) exhausted memory with std::bad_alloc, killing whole batch
    # runs (observed twice, 2026-08-07). Born-digital PDFs don't need OCR;
    # genuinely scanned ones come out under MIN_CHARS, get the `low-text`
    # flag, and go through the targeted OCR/vision finishing layer one file
    # at a time.
    conv = getattr(_tls, "converter", None)
    if conv is None:
        # Recent docling/torch tries to JIT-compile the layout model via
        # torch inductor, which needs a C++ compiler (MSVC `cl` on Windows)
        # and crashes the whole conversion when none exists. Eager mode is
        # fine for inference — force it unless the user overrode it.
        os.environ.setdefault("TORCHDYNAMO_DISABLE", "1")
        from docling.document_converter import DocumentConverter, PdfFormatOption
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions
        if os.environ.get("CONVERT_OCR"):
            conv = DocumentConverter()
        else:
            opts = PdfPipelineOptions(do_ocr=False)
            conv = DocumentConverter(format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=opts)})
        _tls.converter = conv
    return conv


def _docling_convert(src):
    from docling_core.types.doc import ImageRefMode
    res = _docling_converter().convert(src)
    return res.document.export_to_markdown(image_mode=ImageRefMode.PLACEHOLDER)


# --------------------------------------------------------------------------
# Text-layer recovery (poppler)
# --------------------------------------------------------------------------
# docling's layout model classifies free-floating text boxes on design-heavy
# PDFs (infographics, factsheets, strategy one-pagers) as images and drops
# their text silently: the conversion looks well-formed, keeps its headings,
# and is missing most of the document. Measured on this corpus, one 2 MB
# strategy PDF converted to 730 bytes — 16.5% vocabulary recall — with no
# flag raised, because it cleared MIN_CHARS.
#
# The PDF text layer is ground truth for born-digital files and poppler reads
# it in a few hundred milliseconds, so every conversion is checked against it.
# Recovery is structural, not OCR: it only restores text that is already in
# the file. Absent poppler, conversion behaves exactly as before.
_WORD = re.compile(r"[a-z]{4,}")
_pt_cache = {}


def _pdftotext_exe():
    exe = _pt_cache.get("exe", "")
    if exe == "":
        exe = os.environ.get("CONVERT_PDFTOTEXT") or shutil.which("pdftotext")
        if not exe and os.name == "nt":
            found = glob.glob(os.path.join(
                os.environ.get("LOCALAPPDATA", ""), "Microsoft", "WinGet",
                "Packages", "oschwartz10612.Poppler*", "poppler-*",
                "Library", "bin", "pdftotext.exe"))
            exe = found[0] if found else None
        _pt_cache["exe"] = exe
    return exe


def _pdf_text_layer(src):
    """Raw text layer, or None when poppler is unavailable or fails."""
    exe = _pdftotext_exe()
    if not exe:
        return None
    try:
        out = subprocess.run([exe, wp(src), "-"], capture_output=True,
                             timeout=300)
    except (OSError, subprocess.SubprocessError):
        return None
    if out.returncode != 0:
        return None
    return out.stdout.decode("utf-8", errors="replace")


def _vocab(text):
    return set(_WORD.findall(text.lower()))


def _text_layer_blocks(raw):
    """Blank-line-separated blocks, each rewrapped into one paragraph.

    pdftotext hard-wraps at the layout's line breaks; joining within a block
    keeps sentences intact so corpus chunking splits on meaning rather than
    on where a text box happened to end.
    """
    blocks = []
    for chunk in re.split(r"\n\s*\n", raw.replace("\f", "\n\n")):
        para = " ".join(line.strip() for line in chunk.splitlines()
                        if line.strip())
        if para:
            blocks.append(para)
    return blocks


def _recover_text_layer(text, raw):
    """Return (markdown, mode) restoring what docling dropped.

    Below RECALL_REPLACE docling has kept so little that its headings are not
    worth the missing body, and the text layer replaces it wholesale.
    Otherwise the structured conversion stands and only the blocks it missed
    are appended, so tables and headings survive.
    """
    blocks = _text_layer_blocks(raw)
    if not blocks:
        return text, None
    have = _vocab(text)
    gt = _vocab(raw)
    if len(gt & have) / len(gt) < RECALL_REPLACE:
        return "\n\n".join(blocks), "replaced"
    missing = []
    for block in blocks:
        words = _vocab(block)
        if not words:
            continue
        gone = words - have
        if len(gone) >= 3 and len(gone) / len(words) > 0.4:
            missing.append(block)
    if not missing:
        return text, None
    return (text.rstrip() + "\n\n## Text recovered from the PDF text layer\n\n"
            + "\n\n".join(missing) + "\n"), "appended"


def convert_pdf(src):
    text = _docling_convert(src)
    flags = []
    if _visible_len(text) < MIN_CHARS:
        flags.append({
            "type": "low-text",
            "detail": f"{_visible_len(text)} visible chars — likely a scanned/"
                      "graphic PDF; queue for the OCR/vision finishing layer",
        })
    raw = _pdf_text_layer(src)
    gt = _vocab(raw) if raw else set()
    if len(gt) >= 50:
        recall = len(gt & _vocab(text)) / len(gt)
        if recall < RECALL_MIN:
            text, mode = _recover_text_layer(text, raw)
            if mode:
                new = len(gt & _vocab(text)) / len(gt)
                flags.append({
                    "type": "text-recovery",
                    "detail": f"docling recall {recall:.1%} -> {new:.1%} "
                              f"({mode} from the PDF text layer); layout-"
                              "dependent structure in the recovered part is "
                              "flattened to paragraphs",
                })
    return text, flags


_NUM_SUB = re.compile(r"^(\d+(?:\.\d+)+)[.)]?\s")
_NUM_TOP = re.compile(r"^(\d+)[.)]\s")


def _numbering_check(text):
    """Heuristic for docling's fabricated clause numbers on auto-numbered
    Word documents: flag sub-numbers whose top-level component disagrees with
    the most recent top-level section seen. A flag means "review before
    citing clause numbers", not "definitely wrong".
    """
    top = None
    mismatch = total = 0
    for line in text.splitlines():
        line = line.strip()
        if line[:2] in ("- ", "* "):
            line = line[2:]
        m = _NUM_SUB.match(line)
        if m:
            total += 1
            if top is not None and int(m.group(1).split(".")[0]) != top:
                mismatch += 1
            continue
        m = _NUM_TOP.match(line)
        if m:
            top = int(m.group(1))
    if mismatch:
        return [{
            "type": "numbering",
            "detail": f"{mismatch}/{total} hierarchical numbers disagree with "
                      "their enclosing top-level section — docling may have "
                      "fabricated clause numbers; finishing layer before citing",
        }]
    return []


def convert_docx(src):
    text = _docling_convert(src)
    flags = []
    if _visible_len(text) < MIN_CHARS:
        flags.append({
            "type": "low-text",
            "detail": f"{_visible_len(text)} visible chars — scan-in-a-wrapper; "
                      "queue for the OCR/vision finishing layer",
        })
    flags.extend(_numbering_check(text))
    return text, flags


# --------------------------------------------------------------------------
# PPTX via python-pptx
# --------------------------------------------------------------------------
CHART_MAX_ROWS = 60    # beyond this a chart is a dataset, not a slide message


def _num(v):
    if v is None:
        return ""
    if isinstance(v, float) and v == int(v):
        return str(int(v))
    if isinstance(v, float):
        return f"{v:g}"
    return str(v)


def _pptx_chart_md(chart, lines):
    """Charts carry their plot data in the file; render it as a table.

    A slide whose message is one chart used to convert to `<!-- chart -->`,
    which reads as an empty slide and sends the deck to the vision layer for
    numbers that were sitting in the XML all along. Extraction is exact —
    prefer it over describing a picture of the same chart.
    """
    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass
    lines.append("**Chart" + (f": {title}" if title else "") + "**")
    lines.append("")

    try:
        cats = [str(c) for c in chart.plots[0].categories]
    except Exception:
        cats = []
    series = []
    try:
        for s in chart.series:
            series.append((str(s.name or "Series"), list(s.values)))
    except Exception:
        series = []

    if not series:
        lines.append("<!-- chart -->")
        lines.append("")
        return

    rows = max(len(cats), max(len(v) for _, v in series))
    if not cats:
        cats = [str(i + 1) for i in range(rows)]
    cut = rows > CHART_MAX_ROWS
    rows = min(rows, CHART_MAX_ROWS)

    lines.append("| Category | " + " | ".join(n.replace("|", "\\|")
                                              for n, _ in series) + " |")
    lines.append("|---|" + "---|" * len(series))
    for i in range(rows):
        cat = cats[i] if i < len(cats) else ""
        vals = [_num(v[i]) if i < len(v) else "" for _, v in series]
        lines.append("| " + cat.replace("|", "\\|") + " | "
                     + " | ".join(vals) + " |")
    if cut:
        lines.append(f"| … | {' | '.join(['…'] * len(series))} |")
    lines.append("")


def _pptx_shape_md(shape, lines):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        stype = shape.shape_type
    except Exception:
        stype = None
    if stype == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _pptx_shape_md(sub, lines)
        return
    if getattr(shape, "has_table", False):
        rows = [[c.text.strip().replace("\n", " ").replace("|", "\\|")
                 for c in r.cells] for r in shape.table.rows]
        if rows:
            lines.append("| " + " | ".join(rows[0]) + " |")
            lines.append("|" + "---|" * len(rows[0]))
            for r in rows[1:]:
                lines.append("| " + " | ".join(r) + " |")
            lines.append("")
        return
    if getattr(shape, "has_chart", False):
        _pptx_chart_md(shape.chart, lines)
        return
    if stype == MSO_SHAPE_TYPE.PICTURE:
        lines.append("<!-- image -->")
        lines.append("")
        return
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            lines.append(text)
            lines.append("")


def convert_pptx(src):
    from pptx import Presentation
    with open(wp(src), "rb") as fh:
        prs = Presentation(fh)

    out, low = [], []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()
        except Exception:
            pass
        out.append(f"## Slide {i}" + (f" — {title}" if title else ""))
        out.append("")

        lines = []
        for shape in slide.shapes:
            if title and shape is slide.shapes.title:
                continue
            _pptx_shape_md(shape, lines)
        out.extend(lines)

        # Visual density is judged on the slide face only, notes excluded.
        body = "".join(l for l in lines if not l.startswith("<!--"))
        if _visible_len(body) + _visible_len(title) < LOWTEXT_SLIDE:
            low.append(i)

        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                out.append("**Speaker notes:**")
                out.append("")
                out.append("> " + note.replace("\n", "\n> "))
                out.append("")

    flags = []
    if low:
        flags.append({
            "type": "pptx-low-text-slides",
            "detail": f"{len(low)}/{len(prs.slides)} slides under "
                      f"{LOWTEXT_SLIDE} chars: {low} — visual content; queue "
                      "for the per-slide vision finishing layer",
        })
    return "\n".join(out), flags


# --------------------------------------------------------------------------
# HTML — already markup; no docling
# --------------------------------------------------------------------------
def _html_to_md(html):
    try:
        from markdownify import markdownify
        return markdownify(html, heading_style="ATX")
    except ImportError:
        pass
    try:
        import html2text
        h = html2text.HTML2Text()
        h.body_width = 0
        return h.handle(html)
    except ImportError:
        raise RuntimeError(
            "no html->md library in this interpreter: pip install markdownify "
            "(preferred) or html2text, or convert with pandoc")


def convert_html(src):
    with io.open(wp(src), encoding="utf-8", errors="replace") as fh:
        raw = fh.read()
    # Faithful conversion only. Boilerplate-stripping for saved web articles
    # is a per-file judgement call — route those through the defuddle skill.
    return _html_to_md(raw), []


# --------------------------------------------------------------------------
# EML / MSG — parse, don't layout-convert
# --------------------------------------------------------------------------
def _email_frontmatter(fields):
    lines = ["---"]
    for key, val in fields:
        if val:
            val = str(val).replace('"', "'").replace("\n", " ").strip()
            lines.append(f'{key}: "{val}"')
    lines.append("---")
    return "\n".join(lines)


def convert_eml(src):
    import email
    from email import policy
    with open(wp(src), "rb") as fh:
        msg = email.message_from_binary_file(fh, policy=policy.default)

    fm = _email_frontmatter([
        ("from", msg.get("From")), ("to", msg.get("To")),
        ("cc", msg.get("Cc")), ("date", msg.get("Date")),
        ("subject", msg.get("Subject")),
    ])

    body_part = msg.get_body(preferencelist=("plain", "html"))
    body = ""
    if body_part is not None:
        content = body_part.get_content()
        # Emails with a UTF-8 body but no declared charset decode as ASCII
        # with replacement characters; if the raw payload is clean UTF-8,
        # prefer that over the mangled default decode.
        if "�" in content:
            try:
                content = body_part.get_payload(decode=True).decode("utf-8")
            except (UnicodeDecodeError, AttributeError):
                pass
        body = _html_to_md(content) if body_part.get_content_subtype() == "html" else content

    attachments = [p.get_filename() for p in msg.iter_attachments()
                   if p.get_filename()]
    flags = []
    text = fm + "\n\n" + body.strip() + "\n"
    if attachments:
        text += "\n**Attachments (not inlined):** " + ", ".join(attachments) + "\n"
        flags.append({
            "type": "attachments",
            "detail": "route each through its own format pipeline: "
                      + ", ".join(attachments),
        })
    return text, flags


def convert_msg(src):
    try:
        import extract_msg
    except ImportError:
        raise RuntimeError("extract-msg not installed: pip install extract-msg")
    m = extract_msg.Message(wp(src))
    fm = _email_frontmatter([
        ("from", m.sender), ("to", m.to), ("cc", m.cc),
        ("date", m.date), ("subject", m.subject),
    ])
    attachments = [a.longFilename or a.shortFilename for a in m.attachments]
    flags = []
    text = fm + "\n\n" + (m.body or "").strip() + "\n"
    if attachments:
        text += "\n**Attachments (not inlined):** " + ", ".join(attachments) + "\n"
        flags.append({
            "type": "attachments",
            "detail": "route each through its own format pipeline: "
                      + ", ".join(attachments),
        })
    return text, flags


HANDLERS = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".doc": convert_docx,
    ".pptx": convert_pptx,
    ".html": convert_html,
    ".htm": convert_html,
    ".eml": convert_eml,
    ".msg": convert_msg,
}

# Copied into the output tree untouched.
PASSTHROUGH = {".md", ".txt"}

# Deliberately NOT converted — the document-vs-dataset policy: form-template
# workbooks (XLSForms) can be rendered to markdown, data workbooks should be
# catalogued (see the form-catalogue skill) and left in their native format,
# and CSV choice lists left untouched. A spreadsheet flattened to markdown
# is neither a readable document nor a queryable dataset.
POLICY_SKIP = {".xlsx", ".xls", ".xlsm", ".csv"}
